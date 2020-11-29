from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from openpyxl import load_workbook
from openpyxl.styles import Font
from selenium import webdriver
from PIL import Image
from docx import Document
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from shutil import copyfile

import glob
import os
import time
import subprocess
import sys
import datetime
import imaplib
import email
import smtplib
import ssl

prs = Presentation("TT" + '.pptx')

slide_layout = prs.slide_layouts[3]

slide = prs.slides.add_slide(slide_layout)

for i in range(15):
    try:
        titulo = slide.placeholders[i]
        print(i)
    except KeyError:
        pass

titulo.text = "Hola"

prs.save("TT-hecho.pptx")
