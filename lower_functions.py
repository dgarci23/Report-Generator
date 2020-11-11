from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from openpyxl import load_workbook
from openpyxl.styles import Font
from selenium import webdriver
from PIL import Image
from docx import Document
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

import glob
import os
import time
import subprocess
import sys
import datetime
import imaplib
import email
import smtplib
import ssl

#
#
# FUNCIÓN PRINCIPAL: DanielaPPTX
#
#



# ------------------------------------Funciones sobre cargar xlsx---------------------------
# Descripción: Funciones que abren archivo .xlsx y borran referencias previas

'''
Nombre de la función: loadXLSX
Entradas: Preferencias
Salidas: sheet_ranges

Carga el archivo .xlsx y elimina archivo pre-existente con su nombre
'''
def loadXLSX(Preferencias, Empresa):

    wb = load_workbook(filename=Preferencias[1] + '.xlsx', read_only = True, data_only = True)

    Empresa["Mes"] = Preferencias[3]

    sheet_ranges = wb[Empresa["Mes"]]

    return sheet_ranges

# -------------------------------------Funciones sobre Medidas------------------------------
# Descripción: Funciones que regresan diccionarios con las medidas de elementos

'''
Nombre de la función: f_IconoMedidas
Entrada: NombreEmpresa (string)
Salida: IconoMedidas (dict)

Manejo de las medidas de los iconos
'''
def f_IconoMedidas(NombreEmpresa):

    IconoMedidas = {"Top": Inches(0.8), "Left": Inches(8.2), "Height": Inches(0.57)}

# Casos particulares de empresas
    if NombreEmpresa == 'COEL':

        IconoMedidas["Top"] = Inches(1.2)

    elif NombreEmpresa == 'UMECIT':

        IconoMedidas["Top"] = Inches(0.05)
        IconoMedidas["Left"] = Inches(6.3)

    elif NombreEmpresa == 'TAGORE':

        IconoMedidas["Top"] = Inches(0.15)

    elif NombreEmpresa == 'PANAMA PACIFICO':

        IconoMedidas["Top"] = Inches(0.95)


    return IconoMedidas


'''
Nombre de la función: f_TablaMedidas
Entrada: NombreEmpresa
Salida: TablaMedidas (dict)

Manejo de las medidas de la tabla
'''
def f_TablaMedidas(NombreEmpresa):

    TablaMedidas = {"Width" : Inches(3.6),
                    "Height" : Inches(3.94),
                    "Left" : Inches(6),
                    "Top" : Inches(2),
                    "Rows" : 10,
                    "Cols" : 2}

    if NombreEmpresa == 'TAGORE':
        TablaMedidas["Top"] = Inches(1.2)

    return TablaMedidas


'''
Nombre de la función: f_CapturaMedidas
Entrada: NombreEmpresa
Salida: CapturaMedidas

Manejo de las medidas de la captura
'''
def f_CapturaMedidas(NombreEmpresa):

    CapturaMedidas = {"Top" : Inches(2.3),
                      "Left" : Inches(0.25)}

    if NombreEmpresa == 'TAGORE':
        CapturaMedidas["Top"] = Inches(1.6)

    return CapturaMedidas


'''
Nombre de la función: f_FechaMedidas
Entrada: NombreEmpresa
Salida: FechaMedidas

Medidas para la fecha en la diapositiva
'''
def f_FechaMedidas(NombreEmpresa):

    FechaMedidas = {"Left" : Inches(0.4),
                    "Width" : Inches(2),
                    "Height" : Inches(0.3),
                    "Top" : Inches(1.2)}

    #if (NombreEmpresa == 'TAGORE'):
    #    FechaMedidas["Top"] = Inches(1.7)

    return FechaMedidas


# --------------------------------Funciones sobre indices----------------------------------
# Descripción: Obtención de índices (cantidad de entradas en el archivo de Excel)

'''
Nombre de la función: f_FinalIndex
Entrada: sheet_ranges
Salida: index_final (double)

Encuentra la cantidad de entradas en el xlsx
'''
def f_FinalIndex(sheet_ranges):

    index_answers = 9
    while (True):

        check_cell = 'A' + str(index_answers)

        if sheet_ranges[check_cell].value.lower() == 'resultados totales':
            index_final = index_answers - 1
            break
        else:
            index_answers += 1


    return index_answers


# --------------------------------Funciones sobre descarga de Imágenes----------------------
# Descripción: Descarga y compresión de imagénes (capturas de páginas web)

'''
Nombre de la función: f_loopCompress
Entrada: None
Salida: None

Loopea en las publicaciones y llama a la función que comprime fotos
Posibilidad: añadir reporte a un .txt
'''
def f_loopCompress():

    #os.chdir('/home/dgarci23/Fotos')

    for file_name in glob.glob("*.png"):
        f_Compress(file_name)

    #os.chdir('/home/dgarci23')


'''
Nombre de la funcion: f_loopScreenShots
Entrada: NombreEmpresa, index_final, sheet_ranges
Salida: None

Loopea por los links y llama a la funcion que toma screenshots
'''
def f_loopScreenShots(NombreEmpresa, index_final, sheet_ranges):

    for current_index in range(9, index_final):

        link = sheet_ranges['AA' + str(current_index)].value.lower()

        if link != 'no':

            f_ScreenShots(NombreEmpresa, current_index-8, link)
            print('Exito: ' + link)



'''
Nombre de la función: f_ScreenShots
Entrada: NombreEmpresa, current_index, link
Salida: None

Carga pagina web con el link
Toma screenshot de la pantalla
Guarda el resultado
'''
def f_ScreenShots(NombreEmpresa, current_index, link):

    options = webdriver.ChromeOptions()
    options.add_argument('headless')

    delay = 100
    # CHANGE WITH RESPECT TO LINUX VERSION '.EXE'
    with webdriver.Chrome('./chromedriver', options=options) as driver:

        file_name = NombreEmpresa + ' ' + str(current_index) + '.png'

        if os.path.exists(file_name): return 0;

        width = 1000
        height = 1000

        driver.set_window_size(width, height)
        link = str(link)

        try:
            driver.get(link)
        except:
            pass

        try:
            if 'instagram' in link:

                f_IgScreenShots(driver, link)

            else:
                time.sleep(10)
        except:
            pass

        driver.save_screenshot(file_name)

        driver.quit()


'''
Nombre de la función: f_IgScreenShots
Entrada: driver, link (str - no necesario)
Salida: none

Realiza el "log-in" en Ig con los credenciales dados
'''
def f_IgScreenShots(driver, link):

    time.sleep(3)

    driver.find_element_by_name("username").send_keys('danielapp2020')
    driver.find_element_by_name("password").send_keys('3XW]FK"j}dJkFDsz')
    driver.find_element_by_name("password").send_keys(u'\ue007')

    time.sleep(3)

    submit_button = driver.find_elements_by_xpath('//*[@id="react-root"]/section/main/div/div/div/div/button')[0]
    submit_button.click()

    time.sleep(3)


'''
Nombre de la función: f_Compress
Entrada: file_name
Salida: None

Comprime imagenes
'''
def f_Compress(file_name):

    image_name = file_name.split('.')[0]
    image_open = Image.open(file_name)
    image_open = image_open.convert('RGB')
    image_open.save(image_name+'.jpg',optimize=True,quality=85)
    os.remove(file_name)


# ------------------------------------Ordenar los elementos---------------------------------
'''
Nombre de la función: sortIndex
Entradas: index_final, sheet_ranges
Salidas: sort_list

Ordena las entradas dependiendo del tipo de medio
'''
def sortIndex(index_final, sheet_ranges):

    sort_list = [[i, sheet_ranges['G' + str(i)].value.lower()] for i in range(9, index_final)]

    sort_list.sort(key=sortFunction)

    return sort_list


'''
Nombre de la función: sortFunction
Entradas: element
Salidas: 1, 0

(FUNCION QUE) Ordena la lista dependiendo de el tipo de medio
'''
def sortFunction(element):

    medios_digitales = ['twitter', 'instagram', 'facebook']

    medios_tradicionales = ['web', 'impreso', 'radio', 'televisión']

    if element[1] in medios_digitales: return 1
    if element[1] in medios_tradicionales: return 0

# ---------------------------Funciones para añadir contenido a diapositivas-----------------
# Descripción: añadir el contenido de las diapositivas (tabla, titulo, imagen, icono)

'''
Nombre de la función: addSlide
Entrada: prs
Salida: slide

Añade una nueva diapositiva
'''
def addSlide(presentation):

    placeholder = 3

    slide_layout = presentation.slide_layouts[placeholder]

    slide = presentation.slides.add_slide(slide_layout)

    return slide


'''
Nombre de la función: createTable
Entradas: slide, TablaMedidas, sheet_ranges, index_copy
Salidas: table

Crea la tabla, llama la función que la llena
'''
def createTable(slide, TablaMedidas, sheet_ranges, index_copy):

    table = slide.shapes.add_table(TablaMedidas["Rows"],
                             TablaMedidas["Cols"],
                             TablaMedidas["Left"],
                             TablaMedidas["Top"],
                             TablaMedidas["Width"],
                             TablaMedidas["Height"]).table

    table.cell(0, 0).text = 'Medio'
    table.cell(1, 0).text = 'Tipo'
    table.cell(2, 0).text = 'Tamaño'
    table.cell(3, 0).text = 'Valor Publicitario'
    table.cell(4, 0).text = 'Valor Informativo'
    table.cell(5, 0).text = 'Tier'
    table.cell(6, 0).text = 'Mención de Marca'
    table.cell(7, 0).text = 'Presencia de vocero'
    table.cell(8, 0).text = 'Quote de vocero'
    table.cell(9, 0).text = 'Favorabilidad mediática'

    table = fillTable(table, sheet_ranges, index_copy, TablaMedidas)

    return table

'''
Nombre de la función: fillTable
Entrada: table, sheet_ranges, index_copy
Salida: table

Llena la tabla con el contenido de el archivo de Excel
'''
def fillTable(table, sheet_ranges, index_copy, TablaMedidas):

    table.cell(0, 1).text = sheet_ranges['H' + str(index_copy)].value
    table.cell(1, 1).text = sheet_ranges['G' + str(index_copy)].value
    table.cell(2, 1).text = sheet_ranges['I' + str(index_copy)].value

    table.cell(3, 1).text = DineroAjuste(str(sheet_ranges['L' + str(index_copy)].value))
    table.cell(4, 1).text = DineroAjuste(str(sheet_ranges['M' + str(index_copy)].value))

    table.cell(5, 1).text = str(sheet_ranges['N' + str(index_copy)].value)

    table.cell(6,1).text = binaryCheck(sheet_ranges, 'R', index_copy)
    table.cell(7,1).text = binaryCheck(sheet_ranges, 'T', index_copy)
    table.cell(8,1).text = binaryCheck(sheet_ranges, 'V', index_copy)

    if sheet_ranges['X' + str(index_copy)].value == 1:
        table.cell(9, 1).text = 'Positiva'
    elif sheet_ranges['Y' + str(index_copy)].value == 1:
        table.cell(9, 1).text = 'Negativa'
    else:
        table.cell(9, 1).text = 'Neutra'

    table = formatTable(table, TablaMedidas)

'''
Nombre de la función: binaryCheck
Entrada: sheet_ranges, column, index_copy
Salida: string

Revisa una celda y regresa 'Si' o 'No'
'''
def binaryCheck(sheet_ranges, column, index_copy):

    if sheet_ranges[column + str(index_copy)].value == None:
        return 'No'
    else:
        return 'Si'


'''
Nombre de la función: formatTable
Entradas: table, TablaMedidas
Salidas: table

Formatea la tabla
'''
def formatTable(table, TablaMedidas):

    for i in range(TablaMedidas["Rows"]):
        for j in range(TablaMedidas["Cols"]):

            font = table.cell(i, j).text_frame.paragraphs[0].font
            font.size = Pt(12)
            if j == 0:
                table.cell(i, j).text_frame.paragraphs[0].font.bold = True
            elif i == 0:
                table.cell(i, j).text_frame.paragraphs[0].font.bold = True

    return table


'''
Nombre de la función: DineroAjuste
Entradas: dinero
Salidas: dinerofinal

Convierte el formato para dinero con , y .
'''
def DineroAjuste(dinero):

    dinero = dinero.split('.')

    dolares = dinero[0]

    if len(dinero) == 1:

        cents = '.00'


    else:
        if len(dinero[1]) == 1:
            cents = dinero[1] + '0'
        else:
            cents = dinero[1]

        cents = '.' + cents

    if len(dolares) < 4:
        dinerofinal = 'B/. ' + dolares + cents
    else:
        dinerofinal = 'B/. ' + dolares[0:-3] + ',' + dolares[-3:-1] + dolares[len(dolares) - 1] + cents



    return dinerofinal


'''
Nombre de la función: addIcon
Entradas: file_name, IconoMedidas, slide
Salidas: None

Añade el icono
'''
def addIcon(file_name, IconoMedidas, slide):

    file_name=file_name.strip()

    if file_name in ['Noticiero', 'Televisión']:
        file_name = 'TV'

    if file_name == 'Impreso':
        IconoMedidas["Height"] = Inches(0.85)
        IconoMedidas["Top"] -= Inches(0.1)

    file_name_ext = 'Logo/' + file_name + '.png'

    pic = slide.shapes.add_picture(file_name_ext,
                                   IconoMedidas["Left"],
                                   IconoMedidas["Top"],
                                   height = IconoMedidas["Height"])


'''
Nombre de la función: addTitle
Entradas: title_texto, slide
Salidas: None

Añade el titulo a la diapositiva
'''
def addTitle(title_texto, slide):

    title_texto = 'Publicaciones - ' + title_texto

    titulo = slide.placeholders[0]

    titulo.text = title_texto

    titleFormat(titulo)


'''
Nombre de la función: titleFormat
Entradas: titulo
Salidas: None

Formato del titulo
'''
def titleFormat(titulo):

    titulo.text_frame.paragraphs[0].font.bold = True

    font_titulo = titulo.text_frame.paragraphs[0].font

    font_titulo.size = Pt(35)

    titulo.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


'''
Nombre de la función: addDate
Entradas: Fecha, FechaMedidas, slide
Salidas: None

Añade la fecha a la diapositiva
'''
def addDate(Fecha, FechaMedidas, slide):

    Fecha = editFecha(Fecha)

    slide.placeholders[1].left = FechaMedidas["Left"]
    slide.placeholders[1].width = FechaMedidas["Width"]
    slide.placeholders[1].height = FechaMedidas["Height"]
    slide.placeholders[1].top = FechaMedidas["Top"]

    slide.placeholders[1].text = Fecha
    slide.placeholders[1].text_frame.paragraphs[0].font.bold = True
    font_fecha = slide.placeholders[1].text_frame.paragraphs[0].font
    font_fecha.size = Pt(16)


'''
Nombre de la función: editFecha
Entradas: Fecha
Salidas: Fecha

Edita la fecha, regresa la fecha en formato 'legible'
'''
def editFecha(Fecha):

    mes = Fecha[0:3]
    dia = Fecha[4:6]

    if mes == 'Jan':
        mes = 'Enero'
    elif mes == 'Feb':
        mes = 'Febrero'
    elif mes == 'Mar':
        mes = 'Marzo'
    elif mes == 'Apr':
        mes = 'Abril'
    elif mes == 'May':
        mes = 'Mayo'
    elif mes == 'Jun':
        mes = 'Junio'
    elif mes == 'Jul':
        mes = 'Julio'
    elif mes == 'Aug':
        mes = 'Agosto'
    elif mes == 'Sep':
        mes = 'Septiembre'
    elif mes == 'Oct':
        mes = 'Octubre'
    elif mes == 'Nov':
        mes = 'Noviembre'
    elif mes == 'Dic':
        mes = 'Diciembre'

    Fecha = dia + ' de ' + mes

    return Fecha


'''
Nombre de la función: addLink
Entradas: link, slide
Salidas: None

Añade el link
'''
def addLink(link, slide):

    link_placeholder = slide.placeholders[13]

    link_placeholder.text = link



'''
Nombre de la función: addCaptura
Entradas: NombreEmpresa, index_image, slide, CapturaMedidas
Salidas: None

Añade el screenshot de la página web
'''
def addCaptura(NombreEmpresa, index_image, slide, CapturaMedidas):

    filename = NombreEmpresa + ' ' + str(index_image) + '.jpg'

    try:
        #os.chdir('/home/dgarci23/Fotos')
        [CapturaMedidas["Width"], CapturaMedidas["Height"]] = Image.open(filename).size
        #os.chdir('/home/dgarci23/')

        if Inches(4.7) * CapturaMedidas["Width"]/CapturaMedidas["Height"] < Inches(4):
            slide.shapes.add_picture(filename,
                                     CapturaMedidas["Left"],
                                     CapturaMedidas["Top"],
                                     height=Inches(4.5))
        else:
            slide.shapes.add_picture(filename,
                                     CapturaMedidas["Left"],
                                     CapturaMedidas["Top"],
                                     width=Inches(4.2))

        os.remove(filename)

    except OSError:
        pass


#
#
# FUNCIÓN PRINCIPAL: DanielaWord
#
#

'''
Nombre de la función: prf_2_txt
Entrada: doc
Salida: txt

Convierte formato de parrafos de .docx a un string unico txt
'''
def prf_2_txt(doc):

    txt = ''

    for prf in doc.paragraphs:

        txt += ' ' + prf.text

    txt = txt.replace(":/", "^")
    txt = txt.replace(':','')
    txt = txt.replace("^", ":/")

    return txt


'''
Nombre de la función: load_categorias
Entrada: None
Salidas: categorias

Consigue la lista de categorias
NOTA: Lista de categorías eliminadas:
* Quote de Vocer
* Mención de Marca
* Presencia de Vocero
'''
def load_categorias():

    categorias = ["Fecha", "Título de Publicación", "Link",
        "Tipo de Gestión", "Titular de Gestión"]

    return categorias


'''
Nombre de la función: SortCategories
Entradas: txt
Salidas: list_index

Regresa una lista con las entradas que corresponden a cada Gestion
'''
def SortCategories(txt):

    index_main = GetIndexes("Tipo de Gestión", txt)
    index_sub = GetIndexes("Fecha", txt)

    order = [[i - j for i in index_sub] for j in index_main]
    order.append([-1 for i in index_sub])

    list_index = []
    i = 1
    j = 0
    max_i = len(order)
    max_j = len(order[0])

    while (i<max_i and j<max_j):
        bool_check = True
        while (bool_check and j<max_j):
            if (order[i][j]<0):
                list_index.append(i-1)
                j += 1
            else:
                i += 1
                bool_check = False

    return list_index


'''
Nombre de la función: GetIndexes
Entradas: string, txt
Salidas: index_list

Revisa los índices donde aparece string en txt y regresa la lista
'''
def GetIndexes(string, txt):
    index_list = []
    i = 0
    while (True):
        index_list.append(txt.find(string, i))
        i = index_list[-1] + 1
        if (index_list[-1] == -1): break;
    del index_list[-1]

    return index_list




'''
Nombre de la función: f_extractElement_loop
Entradas: categorias, txt
Salidas: Data

Llama a extraer los elementos para cada categoria
'''
def f_extractElement_loop(categorias, txt):

    Data = {}

    for categoria in categorias:
        Data[categoria] = f_extractElement(categoria, categorias, txt)

    #print(Data['Título de Publicación:'])

    return Data


'''
Nombre de la función: addIndex
Entradas: data
Salidas: data

Añade el índice de la entrada a data
'''
def addIndex(data):

    data["Index"] = [i for i in range(1, len(data["Fecha"])+1)]

    return data


'''
Nombre de la función: f_extractElement
Entradas: sel_categoria, categorias, txt
Salidas: elementos_list

Extrae los elementos de una categoria
'''
def f_extractElement(sel_categoria, categorias, txt):

    elementos_lista = []
    index = 0
    end = 1
    while end >= 0:
        index = txt.find(sel_categoria, index)
        end *= index/abs(index)

        indexStart = index + len(sel_categoria)

        indexStop = len(txt)
        for categoria_check in categorias:
            tmp_indexStop = txt.find(categoria_check, indexStart)
            if tmp_indexStop < indexStop and tmp_indexStop != -1:
                indexStop = tmp_indexStop

        elemento_extraido = txt[indexStart:indexStop]

        elemento_extraido = elemento_extraido.strip()

        elementos_lista.append(elemento_extraido)

        index += 1

    del elementos_lista[-1]

    return elementos_lista


'''
Nombre de la función: BinaryUpdate_loop
Entradas: data
Salidas: data

Regresa Llama a BinaryUpdate para categoria necesaria
'''
def BinaryUpdate_loop(data):

    categorias_binary = ['Quote de Vocero', 'Presencia de Vocero', 'Mención de Marca']

    for categoria in list(data.keys()):

        if categoria in categorias_binary:

            BinaryUpdate(data, categoria)

    return data


'''
Nombre de la función: BinaryUpdate
Entradas: data, categoria
Salidas: data

Cambia 'si' y 'no' por 1 y 0
'''
def BinaryUpdate(data, categoria):

    for index_binary in range(len(data[categoria])):

        elemento = data[categoria][index_binary]

        if elemento.lower() in ['si', 'sí']:

            data[categoria][index_binary] = 1

        elif elemento.lower() in ['no']:

            data[categoria][index_binary] = 0

    return data


'''
Nombre de la función: f_date2number_loop
Entradas: data
Salidas: data

Procesa la fecha para poder llamada a f_date2number
'''
def f_date2number_loop(data):

    for index_date in range(len(data['Fecha'])):

        elemento = data['Fecha'][index_date]

        elemento = f_date2number(elemento)

        data['Fecha'][index_date] = elemento

    return data


'''
Nombre de la función: f_date2number
Entradas: fecha
Salida: fecha

Cambia la fecha en texto por fecha en numero (formato datetime)
'''
def f_date2number(fecha):

    tmp = fecha
    fecha = fecha.split()

    Año = '2020'

    if type(fecha)==list and len(fecha)>2:
        Mes = fecha[2]

        if Mes == 'enero':
            Mes = '01'
        elif Mes == 'febrero':
            Mes = '02'
        elif Mes == 'marzo':
            Mes = '03'
        elif Mes == 'abril':
            Mes = '04'
        elif Mes == 'mayo':
            Mes = '05'
        elif Mes == 'junio':
            Mes = '06'
        elif Mes == 'julio':
            Mes = '07'
        elif Mes == 'agosto':
            Mes = '08'
        elif Mes == 'septiembre':
            Mes = '09'
        elif Mes == 'octubre':
            Mes = '10'
        elif Mes == 'noviembre':
            Mes = '11'
        else:
            Mes = '12'

        Dia = fecha[0]

        fecha = datetime.datetime.strptime('2020' + Mes + Dia, '%Y%m%d')
    else:
        fecha = tmp

    return fecha


'''
Nombre de la función: FillExcel
Entradas: data, categorias, sheet_ranges
Salidas: None

Llena el documento de excel usando la información en data
'''
def FillExcel(data, categorias, sheet_ranges, orderGestion):

    categorias.append("Index")
    categorias.remove("Tipo de Gestión")
    categorias.remove("Titular de Gestión")

    col_cat = {"Fecha" : ['F'],
                        "Título de Publicación" : ['J'],
                        "Link" : ['AA'],
                        "Tipo de Gestión" : ['D'],
                        "Titular de Gestión" : ['E'],
                        "Nombre de Medio" : ['H'],
                        "Index" : ['K']}

    for categoria in categorias:

        for index_elemento in range(len(data[categoria])):

            cell_id = col_cat[categoria][1*(len(col_cat[categoria])==2)] + str(index_elemento+9)

            sheet_ranges[cell_id].value = data[categoria][index_elemento]

    OperationsExcel(data, sheet_ranges)

    Gestion_loop(data, sheet_ranges, orderGestion)


'''
Nombre de la función: Operations Excel
Entradas; data, sheet_ranges
Salidas: None

Realiza operaciones adentro de excel y asigna formulas
'''
def OperationsExcel(data, sheet_ranges):

    for index_elemento in range(len(data["Link"])):

        addHyperlink(sheet_ranges, index_elemento, data["Link"][index_elemento])

        TipoMedio = GetTipoMedio(data["Link"][index_elemento])

        sheet_ranges['G' + str(index_elemento+9)].value = TipoMedio

        sheet_ranges['H' + str(index_elemento+9)].value = GetNombreMedio(data["Link"])

        sheet_ranges['I' + str(index_elemento+9)].value = GetSize(TipoMedio)

        sheet_ranges['L' + str(index_elemento+9)].value = GetValor(TipoMedio, index_elemento)

        sheet_ranges['O' + str(index_elemento+9)].value = '=IF(N{}=1,1,"")'.format(index_elemento+9)

        sheet_ranges['Q' + str(index_elemento+9)].value = '=IF(N{}=3,1,"")'.format(index_elemento+9)

        sheet_ranges['M' + str(index_elemento+9)].value = '=L{}*4'.format(index_elemento+9)

'''
Nombre de la función: addHyperlink
Entradas: sheet_ranges, index_elemento, link
Salidias: None

Añade el hipervinculo a la columna G
'''
def addHyperlink(sheet_ranges, index, link):
    if link.lower().strip() != 'no':
        sheet_ranges['G' + str(index+9)].hyperlink = link
        sheet_ranges['G' + str(index+9)].font = Font(size = "11", color="0563C1")


'''
Nombre de la función: GetTipoMedio
Entradas: link
Salidas: TipoMedio

Regresa el Tipo de Medio dependiendo del link
'''
def GetTipoMedio(link):

    if 'instagram' in link.lower():
        return "Instagram"
    elif 'twitter' in link.lower():
        return "Twitter"
    elif "facebook" in link.lower():
        return "Facebook"
    elif "http" in link.lower():
        return "Web"
    else:
        return ""


'''
Nombre de la función: GetSize
Entradas: TipoMedio
Salidas: Tamaño

Regresa el Tamaño dependiendo del Tipo de Medio
'''
def GetSize(TipoMedio):

    if TipoMedio.lower().strip() == "web":
        return "1 página"
    elif TipoMedio.lower().strip() in ["instagram","facebook","twitter"]:
        return "1 post"
    else:
        return ""


'''
Nombre de la función: GetSize
Entradas: TipoMedio
Salidas: Tamaño

Regresa el Tamaño dependiendo del Tipo de Medio
'''
def GetNombreMedio(link):

    file = open("NombreMedios.txt", "r", encoding="utf8")

    lines = file.read()

    for line in lines:

        for medio in line[1:len(lines)]:

            if medio in link:

                return line[0]

    return ""


'''
Nombre de la función: GetSize
Entradas: TipoMedio, index_elemento
Salidas: Tamaño

Regresa el Tamaño dependiendo del Tipo de Medio
'''
def GetValor(TipoMedio, index_elemento):

    if TipoMedio.lower().strip() == "web":
        return '=IF(N{}=3, 750, IF(N{}=1, 1000, ""))'.format(index_elemento+9,index_elemento+9)
    elif TipoMedio.lower().strip() in ["twitter", "instagram", "facebook"]:
        return 500
    else:
        return ""


'''
Nombre de la función: Gestion_loop
Entradas: data, sheet_ranges, orderGestion
Salidas: None

Agrega el Titular de Gestion y el Tipo de Gestion
'''
def Gestion_loop(data, sheet_ranges, orderGestion):
    gestion_cat = [["Titular de Gestión", 'E'],["Tipo de Gestión", 'D']]

    for gestion_tipo in gestion_cat:

        for index in range(len(data["Fecha"])):

            sheet_ranges[gestion_tipo[1] + str(index + 9)].value = data[gestion_tipo[0]][orderGestion[index]]


#
#
# FUNCIÓN PRINCIPAL: Download_Email
#
#

'''
Nombre de la función: f_Email_Preferencias
Entradas: None
Salidas: mail

Carga las preferencias de la funcion principal
'''
def f_Email_Preferencias():

    mail = imaplib.IMAP4_SSL('imap.gmail.com')
    mail.login("DanielaPPTX@gmail.com","1234567Dg!")
    mail.select("INBOX")

    return mail


'''
Nombre de la función: get_messages
Entradas: mail
Salidas: msgs

Obtiene los mensajes
'''
def get_messages(mail):

    msgs = mail.search(None, 'ALL')[1]

    msgs = msgs[0].split()

    return msgs


'''
Nombre de la función: get_attachments
Entradas: mail, msgs, email_id
Salidas: Preferencias

Downloads the attachments
'''
def get_attachments(mail, msgs, email_id):

    data = mail.fetch(email_id, "(RFC822)")[1]

    email_body = data[0][1]

    m = email.message_from_bytes(email_body)

    Preferencias = m['Subject']

    if not(m.get_content_maintype() != 'multipart'):
        #continue

        for part in m.walk():

            if part.get_content_maintype() == 'multipart':
                continue

            if part.get('Content-Disposition') is None:
                continue

            filename = part.get_filename()

            if filename is not None:
                sv_path = os.path.join(filename)
                if not os.path.isfile(sv_path):
                    fp = open(sv_path, 'wb')
                    fp.write(part.get_payload(decode=True))
                    fp.close()

    return Preferencias


#
#
# FUNCIÓN PRINCIPAL: Send_Email
#
#

'''
Nombre de la función: get_subject
Entradas: file_size, tipo
Salidas: subject

Regresa el subject dependiendo del tamaño del archivo y el nombre dependiendo del tipo
'''
def f_size_type(file_size, tipo):

    if file_size < 25:
        subject = "Respuesta procesada"
    else:
        subject = "Archivo muy pesado, contactar"


    if tipo == 'Excel':
        filename = "Presentacion Terminada.pptx"
    elif tipo == 'Word':
        filename = 'Excel Terminado.xlsx'

    return subject, filename


'''
Nombre de la función: Credenciales_Correo
Entradas: None
Salidas: sender, receiver, pwrd

Carga las credenciales del correo
'''
def Credenciales_Correo():

    Cred = {"sender_email" : "DanielaPPTX@gmail.com",
            "receiver_email" : "daniacostav09@gmail.com",
            "pwrd" : "1234567Dg!"}

    return Cred


'''
Nombre de la función: CreateMessage()
Entradas: file_size, tipo
Salidas: message, filename

Crea el mensaje
'''
def CreateMessage(file_size, tipo, Cred):

    message = MIMEMultipart()

    message["From"] = Cred["sender_email"]

    message["To"] = Cred["receiver_email"]

    [message["Subject"], filename] = f_size_type(file_size, tipo)

    text = MIMEText(CuerpoCorreo(),"plain")
    message.attach(text)

    return message, filename


'''
Nombre de la función: CuerpoCorreo
Entradas: none
Salidas: cuerpo

Crea el cuerpo del correo
'''
def CuerpoCorreo():

    wb = load_workbook("Excel Terminado.xlsx")
    sheet_ranges = wb.active

    index = 9
    titulos_publicacion = []
    cuerpo = ""

    while (True):

        if sheet_ranges['J' + str(index)].value is None:

            break

        if sheet_ranges['J' + str(index)].value not in titulos_publicacion:

            titulos_publicacion.append(sheet_ranges['J' + str(index)].value)

        index += 1

    cuerpo = "\n".join(titulos_publicacion)

    cuerpo = 'Títulos de Publicación: ' + '\n' + cuerpo

    return cuerpo


'''
Nombre de la función: addAttachment
Entradas: message, filename
Salidas: part

Añade el archivo adjunto
'''
def addAttachment(message, filename):

    with open(filename, "rb") as attachment:

        part = MIMEBase("application", "octet-stream")
        part.set_payload(attachment.read())

    return part


'''
Nombre de la función: encodeMsg
Entradas: part, msg
Salidas: txt

Codifica el mensaje
'''
def encodeMsg(part, msg):

    encoders.encode_base64(part)

    msg.attach(part)

    txt = msg.as_string()

    return txt


'''
Nombre de la función: SecureSend
Entradas: txt, Cred
Salidas: None

Log in al server y envia el correo
'''
def SecureSend(txt, Cred):

    context = ssl.create_default_context()

    with smtplib.SMTP_SSL("smtp.gmail.com", 465, context = context) as server:

        server.login(Cred["sender_email"], Cred["pwrd"])

        server.sendmail(Cred["sender_email"], Cred["receiver_email"], txt)
